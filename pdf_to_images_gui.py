from __future__ import annotations

import shutil
import subprocess
import sys
import tempfile
import threading
import tkinter as tk
from datetime import datetime
from pathlib import Path
from tkinter import filedialog, messagebox, ttk


DEFAULT_DPI = 200
IMAGE_DPI = 200
DOC_DPI = 160
PPT_DPI = 160
APP_ICON_RESOURCE = "assets/pdf_tools_icon.ico"

LANGUAGE_LABELS = {"zh": "中文", "en": "English"}
LABEL_TO_LANGUAGE = {label: code for code, label in LANGUAGE_LABELS.items()}

TEXT = {
    "zh": {
        "app_title": "PDF 工具箱",
        "window_title": "PDF 工具箱",
        "language": "语言",
        "tab_images": "PDF 转图片",
        "tab_merge": "合并 PDF",
        "tab_split": "拆分 PDF",
        "tab_pdf_to_office": "PDF 转 Word/PPT",
        "tab_office_to_pdf": "Word/PPT 转 PDF",
        "tab_decrypt": "解密 PDF",
        "pdf_file": "PDF 文件",
        "pdf_files": "PDF 文件",
        "pdf_files_to_merge": "待合并 PDF",
        "office_file": "Word/PPT 文件",
        "output_folder": "输出文件夹",
        "merge_output_folder": "输出路径",
        "merge_file_name": "合并文件名",
        "merge_name_hint": "留空 = 合并PDF_日期时间.pdf",
        "word_file_name": "Word 文件名",
        "word_name_hint": "留空 = PDF转Word_日期时间.docx",
        "ppt_file_name": "PPT 文件名",
        "ppt_name_hint": "留空 = PDF转PPT_日期时间.pptx",
        "pdf_file_name": "PDF 文件名",
        "office_pdf_name_hint": "留空 = Word转PDF/PPT转PDF_日期时间.pdf",
        "decrypt_pdf_name_hint": "留空 = 解密PDF_日期时间.pdf",
        "output_pdf": "输出 PDF",
        "output_word": "输出 Word",
        "output_ppt": "输出 PPT",
        "password": "PDF 密码",
        "password_hint": "加密 PDF 需要输入密码；未加密可留空",
        "dpi": "清晰度 DPI",
        "pages": "页码",
        "pages_hint": "留空 = 全部，例如 1,3,5-8",
        "split_mode": "拆分方式",
        "split_each_page": "每页一个 PDF",
        "split_one_file": "选中页合成一个 PDF",
        "browse": "选择",
        "add_files": "添加文件",
        "remove_selected": "移除选中",
        "clear": "清空",
        "move_up": "上移",
        "move_down": "下移",
        "export_images": "导出图片",
        "merge": "合并",
        "split": "拆分",
        "to_word": "转 Word",
        "to_ppt": "转 PPT",
        "to_pdf": "转 PDF",
        "decrypt": "解密",
        "initial_status": "请选择功能并填写文件路径。",
        "running": "正在处理...",
        "done": "完成：{detail}",
        "failed": "处理失败。",
        "select_pdf_title": "选择 PDF",
        "select_pdfs_title": "选择多个 PDF",
        "select_output_title": "选择输出文件夹",
        "select_output_pdf_title": "保存 PDF",
        "select_output_docx_title": "保存 Word",
        "select_output_pptx_title": "保存 PPT",
        "select_office_title": "选择 Word/PPT 文件",
        "all_files": "所有文件",
        "office_files": "Word/PPT 文件",
        "word_files": "Word 文件",
        "ppt_files": "PPT 文件",
        "valid_pdf": "请选择有效的 PDF 文件。",
        "valid_office": "请选择有效的 Word/PPT 文件。",
        "valid_output_folder": "请选择有效的输出文件夹。",
        "valid_output_file": "请选择输出文件。",
        "invalid_file_name": "文件名不能包含以下字符：\\ / : * ? \" < > |",
        "need_two_pdfs": "请至少添加两个 PDF 文件。",
        "dpi_number": "DPI 必须是数字。",
        "dpi_range": "DPI 必须在 72 到 600 之间。",
        "empty_page_item": "页码选择中包含空项。",
        "invalid_page_range": "页码范围格式错误：{value}",
        "range_low_to_high": "页码范围必须从小到大：{value}",
        "invalid_page_number": "页码格式错误：{value}",
        "page_outside_range": "页码 {page} 超出范围 1-{page_count}。",
        "empty_pdf": "这个 PDF 没有页面。",
        "encrypted_password_required": "PDF 已加密，请输入密码。",
        "wrong_password": "PDF 密码错误或无法解密。",
        "office_converter_missing": "未找到可用的 Office 转换器。请安装 Microsoft Office 或 LibreOffice。",
        "unsupported_office": "只支持 .doc、.docx、.ppt、.pptx 文件。",
        "unsupported_pdf_to_office": "只支持输出 .docx 或 .pptx。",
        "missing_dependency": "缺少依赖：{name}。请先运行：pip install -r requirements.txt",
        "images_done": "已导出 {count} 张图片到 {output}",
        "merge_done": "已合并 {count} 个 PDF 到 {output}",
        "split_done": "已生成 {count} 个 PDF 文件到 {output}",
        "office_done": "已生成 {output}",
        "decrypt_done": "已解密到 {output}",
        "progress_page": "正在处理第 {page} 页（{current}/{total}）。",
    },
    "en": {
        "app_title": "PDF Toolkit",
        "window_title": "PDF Toolkit",
        "language": "Language",
        "tab_images": "PDF to Images",
        "tab_merge": "Merge PDF",
        "tab_split": "Split PDF",
        "tab_pdf_to_office": "PDF to Word/PPT",
        "tab_office_to_pdf": "Word/PPT to PDF",
        "tab_decrypt": "Decrypt PDF",
        "pdf_file": "PDF file",
        "pdf_files": "PDF files",
        "pdf_files_to_merge": "PDF files to merge",
        "office_file": "Word/PPT file",
        "output_folder": "Output folder",
        "merge_output_folder": "Output path",
        "merge_file_name": "Merged file name",
        "merge_name_hint": "Blank = MergedPDF_date_time.pdf",
        "word_file_name": "Word file name",
        "word_name_hint": "Blank = PDFToWord_date_time.docx",
        "ppt_file_name": "PPT file name",
        "ppt_name_hint": "Blank = PDFToPPT_date_time.pptx",
        "pdf_file_name": "PDF file name",
        "office_pdf_name_hint": "Blank = WordToPDF/PPTToPDF_date_time.pdf",
        "decrypt_pdf_name_hint": "Blank = DecryptedPDF_date_time.pdf",
        "output_pdf": "Output PDF",
        "output_word": "Output Word",
        "output_ppt": "Output PPT",
        "password": "PDF password",
        "password_hint": "Enter password for encrypted PDFs; leave blank otherwise",
        "dpi": "DPI",
        "pages": "Pages",
        "pages_hint": "Blank = all, e.g. 1,3,5-8",
        "split_mode": "Split mode",
        "split_each_page": "One PDF per page",
        "split_one_file": "Selected pages into one PDF",
        "browse": "Browse",
        "add_files": "Add Files",
        "remove_selected": "Remove Selected",
        "clear": "Clear",
        "move_up": "Move Up",
        "move_down": "Move Down",
        "export_images": "Export Images",
        "merge": "Merge",
        "split": "Split",
        "to_word": "To Word",
        "to_ppt": "To PPT",
        "to_pdf": "To PDF",
        "decrypt": "Decrypt",
        "initial_status": "Choose a tool and fill in the file paths.",
        "running": "Processing...",
        "done": "Done: {detail}",
        "failed": "Processing failed.",
        "select_pdf_title": "Select PDF",
        "select_pdfs_title": "Select PDF files",
        "select_output_title": "Select output folder",
        "select_output_pdf_title": "Save PDF",
        "select_output_docx_title": "Save Word",
        "select_output_pptx_title": "Save PPT",
        "select_office_title": "Select Word/PPT file",
        "all_files": "All files",
        "office_files": "Word/PPT files",
        "word_files": "Word files",
        "ppt_files": "PPT files",
        "valid_pdf": "Please select a valid PDF file.",
        "valid_office": "Please select a valid Word/PPT file.",
        "valid_output_folder": "Please select a valid output folder.",
        "valid_output_file": "Please select an output file.",
        "invalid_file_name": "File name cannot contain: \\ / : * ? \" < > |",
        "need_two_pdfs": "Please add at least two PDF files.",
        "dpi_number": "DPI must be a number.",
        "dpi_range": "DPI must be between 72 and 600.",
        "empty_page_item": "Page selection contains an empty item.",
        "invalid_page_range": "Invalid page range: {value}",
        "range_low_to_high": "Page range must go from low to high: {value}",
        "invalid_page_number": "Invalid page number: {value}",
        "page_outside_range": "Page {page} is outside 1-{page_count}.",
        "empty_pdf": "This PDF has no pages.",
        "encrypted_password_required": "This PDF is encrypted. Please enter the password.",
        "wrong_password": "The PDF password is wrong or decryption failed.",
        "office_converter_missing": "No Office converter found. Install Microsoft Office or LibreOffice.",
        "unsupported_office": "Only .doc, .docx, .ppt, and .pptx files are supported.",
        "unsupported_pdf_to_office": "Only .docx or .pptx output is supported.",
        "missing_dependency": "Missing dependency: {name}. Run: pip install -r requirements.txt",
        "images_done": "Exported {count} images to {output}",
        "merge_done": "Merged {count} PDFs into {output}",
        "split_done": "Created {count} PDF files in {output}",
        "office_done": "Created {output}",
        "decrypt_done": "Decrypted to {output}",
        "progress_page": "Processing page {page} ({current}/{total}).",
    },
}


class AppError(ValueError):
    def __init__(self, message_key: str, **params: object) -> None:
        super().__init__(message_key)
        self.message_key = message_key
        self.params = params


def translate(language: str, key: str, **params: object) -> str:
    return TEXT[language][key].format(**params)


def resource_path(relative_path: str) -> Path:
    base_path = Path(getattr(sys, "_MEIPASS", Path(__file__).resolve().parent))
    return base_path / relative_path


def require_import(module_name: str, package_name: str | None = None):
    try:
        return __import__(module_name)
    except ImportError as exc:
        raise AppError("missing_dependency", name=package_name or module_name) from exc


def parse_page_selection(page_spec: str, page_count: int) -> list[int]:
    spec = page_spec.strip()
    if not spec:
        return list(range(page_count))

    selected_pages: list[int] = []
    seen: set[int] = set()

    for raw_part in spec.split(","):
        part = raw_part.strip()
        if not part:
            raise AppError("empty_page_item")

        if "-" in part:
            range_parts = [item.strip() for item in part.split("-", 1)]
            if len(range_parts) != 2 or not range_parts[0] or not range_parts[1]:
                raise AppError("invalid_page_range", value=part)

            try:
                start = int(range_parts[0])
                end = int(range_parts[1])
            except ValueError as exc:
                raise AppError("invalid_page_range", value=part) from exc

            if start > end:
                raise AppError("range_low_to_high", value=part)

            page_numbers = range(start, end + 1)
        else:
            try:
                page = int(part)
            except ValueError as exc:
                raise AppError("invalid_page_number", value=part) from exc

            page_numbers = range(page, page + 1)

        for page_number in page_numbers:
            if page_number < 1 or page_number > page_count:
                raise AppError("page_outside_range", page=page_number, page_count=page_count)

            page_index = page_number - 1
            if page_index not in seen:
                seen.add(page_index)
                selected_pages.append(page_index)

    return selected_pages


def open_pdf_document(pdf: Path, password: str = ""):
    fitz = require_import("fitz", "PyMuPDF")
    document = fitz.open(pdf)
    if document.needs_pass:
        if not password:
            document.close()
            raise AppError("encrypted_password_required")
        if not document.authenticate(password):
            document.close()
            raise AppError("wrong_password")
    return document


def open_pdf_reader(pdf: Path, password: str = ""):
    pypdf = require_import("pypdf")
    reader = pypdf.PdfReader(str(pdf))
    if reader.is_encrypted:
        if not password:
            raise AppError("encrypted_password_required")
        try:
            result = reader.decrypt(password)
        except Exception as exc:
            raise AppError("wrong_password") from exc
        if result == 0:
            raise AppError("wrong_password")
    if len(reader.pages) == 0:
        raise AppError("empty_pdf")
    return reader


def export_pdf_pages_to_images(
    pdf: Path,
    output: Path,
    dpi: int,
    page_spec: str,
    password: str = "",
    progress_callback=None,
) -> int:
    document = open_pdf_document(pdf, password)
    try:
        page_count = document.page_count
        if page_count == 0:
            raise AppError("empty_pdf")

        selected_pages = parse_page_selection(page_spec, page_count)
        fitz = require_import("fitz", "PyMuPDF")
        matrix = fitz.Matrix(dpi / 72, dpi / 72)

        for export_index, page_index in enumerate(selected_pages, start=1):
            page = document.load_page(page_index)
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            image_path = output / f"{page_index + 1}.png"
            pixmap.save(image_path)
            report_progress(progress_callback, page_index + 1, export_index, len(selected_pages))

        return len(selected_pages)
    finally:
        document.close()


def merge_pdfs(pdf_paths: list[Path], output_pdf: Path, password: str = "") -> int:
    pypdf = require_import("pypdf")
    writer = pypdf.PdfWriter()

    for pdf_path in pdf_paths:
        reader = open_pdf_reader(pdf_path, password)
        for page in reader.pages:
            writer.add_page(page)

    with output_pdf.open("wb") as output_file:
        writer.write(output_file)

    return len(pdf_paths)


def split_pdf(
    pdf: Path,
    output_dir: Path,
    page_spec: str = "",
    password: str = "",
    one_file: bool = False,
) -> int:
    pypdf = require_import("pypdf")
    reader = open_pdf_reader(pdf, password)
    selected_pages = parse_page_selection(page_spec, len(reader.pages))

    if one_file:
        writer = pypdf.PdfWriter()
        for page_index in selected_pages:
            writer.add_page(reader.pages[page_index])
        output_file = output_dir / f"{pdf.stem}_selected_pages.pdf"
        with output_file.open("wb") as handle:
            writer.write(handle)
        return 1

    for page_index in selected_pages:
        writer = pypdf.PdfWriter()
        writer.add_page(reader.pages[page_index])
        output_file = output_dir / f"{pdf.stem}_page_{page_index + 1}.pdf"
        with output_file.open("wb") as handle:
            writer.write(handle)

    return len(selected_pages)


def decrypt_pdf(pdf: Path, output_pdf: Path, password: str) -> Path:
    pypdf = require_import("pypdf")
    reader = open_pdf_reader(pdf, password)
    writer = pypdf.PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    with output_pdf.open("wb") as handle:
        writer.write(handle)
    return output_pdf


def render_pdf_pages_to_pngs(
    pdf: Path,
    output_dir: Path,
    dpi: int,
    page_spec: str = "",
    password: str = "",
    progress_callback=None,
) -> list[Path]:
    document = open_pdf_document(pdf, password)
    try:
        selected_pages = parse_page_selection(page_spec, document.page_count)
        fitz = require_import("fitz", "PyMuPDF")
        matrix = fitz.Matrix(dpi / 72, dpi / 72)
        images: list[Path] = []

        for export_index, page_index in enumerate(selected_pages, start=1):
            page = document.load_page(page_index)
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            image_path = output_dir / f"page_{page_index + 1}.png"
            pixmap.save(image_path)
            images.append(image_path)
            report_progress(progress_callback, page_index + 1, export_index, len(selected_pages))

        return images
    finally:
        document.close()


def pdf_to_word(
    pdf: Path,
    output_docx: Path,
    page_spec: str = "",
    password: str = "",
    progress_callback=None,
) -> Path:
    docx_module = require_import("docx", "python-docx")
    from docx.shared import Inches

    with tempfile.TemporaryDirectory() as temp_dir:
        image_dir = Path(temp_dir)
        images = render_pdf_pages_to_pngs(pdf, image_dir, DOC_DPI, page_spec, password, progress_callback)
        document = docx_module.Document()

        section = document.sections[0]
        section.top_margin = Inches(0.25)
        section.bottom_margin = Inches(0.25)
        section.left_margin = Inches(0.25)
        section.right_margin = Inches(0.25)

        for index, image_path in enumerate(images):
            document.add_picture(str(image_path), width=Inches(7.5))
            if index < len(images) - 1:
                document.add_page_break()

        document.save(output_docx)

    return output_docx


def pdf_to_ppt(
    pdf: Path,
    output_pptx: Path,
    page_spec: str = "",
    password: str = "",
    progress_callback=None,
) -> Path:
    pptx_module = require_import("pptx", "python-pptx")
    from pptx.util import Inches

    with tempfile.TemporaryDirectory() as temp_dir:
        image_dir = Path(temp_dir)
        images = render_pdf_pages_to_pngs(pdf, image_dir, PPT_DPI, page_spec, password, progress_callback)
        presentation = pptx_module.Presentation()
        presentation.slide_width = Inches(10)
        presentation.slide_height = Inches(7.5)
        blank_layout = presentation.slide_layouts[6]

        for image_path in images:
            slide = presentation.slides.add_slide(blank_layout)
            slide.shapes.add_picture(str(image_path), 0, 0, width=presentation.slide_width, height=presentation.slide_height)

        presentation.save(output_pptx)

    return output_pptx


def office_to_pdf(office_file: Path, output_pdf: Path) -> Path:
    suffix = office_file.suffix.lower()
    if suffix not in {".doc", ".docx", ".ppt", ".pptx"}:
        raise AppError("unsupported_office")

    converter = shutil.which("soffice") or shutil.which("libreoffice")
    if converter:
        output_dir = output_pdf.parent
        subprocess.run(
            [
                converter,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(output_dir),
                str(office_file),
            ],
            check=True,
            capture_output=True,
            text=True,
        )
        generated = output_dir / f"{office_file.stem}.pdf"
        if generated.resolve() != output_pdf.resolve():
            if output_pdf.exists():
                output_pdf.unlink()
            generated.replace(output_pdf)
        return output_pdf

    try:
        import win32com.client  # type: ignore[import-not-found]
    except ImportError as exc:
        raise AppError("office_converter_missing") from exc

    try:
        if suffix in {".doc", ".docx"}:
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            document = None
            try:
                document = word.Documents.Open(str(office_file.resolve()))
                document.ExportAsFixedFormat(str(output_pdf.resolve()), 17)
            finally:
                if document is not None:
                    document.Close(False)
                word.Quit()
        else:
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            presentation = None
            try:
                presentation = powerpoint.Presentations.Open(str(office_file.resolve()), WithWindow=False)
                presentation.SaveAs(str(output_pdf.resolve()), 32)
            finally:
                if presentation is not None:
                    presentation.Close()
                powerpoint.Quit()
    except Exception as exc:
        raise AppError("office_converter_missing") from exc

    return output_pdf


def report_progress(progress_callback, page: int, current: int, total: int) -> None:
    if progress_callback:
        percent = int((current / total) * 100)
        progress_callback(percent, page, current, total)


def ensure_pdf(path: Path) -> None:
    if not path.is_file() or path.suffix.lower() != ".pdf":
        raise AppError("valid_pdf")


def ensure_output_folder(path: Path) -> None:
    if not path.exists():
        path.mkdir(parents=True, exist_ok=True)
    if not path.is_dir():
        raise AppError("valid_output_folder")


def ensure_output_file(path: Path, suffix: str) -> Path:
    if not str(path):
        raise AppError("valid_output_file")
    if path.suffix.lower() != suffix:
        path = path.with_suffix(suffix)
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


def default_output_file_name(prefix: str, suffix: str) -> str:
    return f"{prefix}_{datetime.now():%Y%m%d_%H%M}{suffix}"


def default_merge_file_name() -> str:
    return default_output_file_name("合并PDF", ".pdf")


def normalize_output_file_name(name: str, default_name: str, suffix: str) -> str:
    clean_name = name.strip()
    if not clean_name:
        clean_name = default_name

    invalid_chars = set('\\/:*?"<>|')
    if any(char in invalid_chars for char in clean_name):
        raise AppError("invalid_file_name")

    if Path(clean_name).name != clean_name:
        raise AppError("invalid_file_name")

    if Path(clean_name).suffix.lower() != suffix:
        clean_name = f"{clean_name}{suffix}"

    return clean_name


def normalize_pdf_file_name(name: str) -> str:
    return normalize_output_file_name(name, default_merge_file_name(), ".pdf")


def build_merge_output_pdf(pdf_paths: list[Path], output_folder: str, file_name: str) -> Path:
    if output_folder.strip():
        folder = Path(output_folder.strip())
    else:
        folder = pdf_paths[0].parent

    ensure_output_folder(folder)
    return folder / normalize_pdf_file_name(file_name)


def build_named_output_file(
    source: Path,
    output_folder: str,
    file_name: str,
    default_prefix: str,
    suffix: str,
) -> Path:
    folder = Path(output_folder.strip()) if output_folder.strip() else source.parent
    ensure_output_folder(folder)
    default_name = default_output_file_name(default_prefix, suffix)
    return folder / normalize_output_file_name(file_name, default_name, suffix)


def office_to_pdf_default_prefix(source: Path) -> str:
    if source.suffix.lower() in {".doc", ".docx"}:
        return "Word转PDF"
    return "PPT转PDF"


class PdfToolsApp(tk.Tk):
    def __init__(self) -> None:
        super().__init__()
        self.language = "zh"
        self.title(self.tr("app_title"))
        self.set_window_icon()
        self.geometry("920x650")
        self.minsize(860, 600)

        self.text_widgets: dict[str, list[tk.Widget]] = {}
        self.tabs: dict[str, ttk.Frame] = {}
        self.status_key = "initial_status"
        self.status_params: dict[str, object] = {}
        self.status = tk.StringVar(value=self.tr("initial_status"))
        self.progress = tk.IntVar(value=0)
        self.language_label = tk.StringVar(value=LANGUAGE_LABELS[self.language])

        self.image_pdf = tk.StringVar()
        self.image_output = tk.StringVar()
        self.image_dpi = tk.IntVar(value=IMAGE_DPI)
        self.image_pages = tk.StringVar()
        self.image_password = tk.StringVar()

        self.merge_files: list[Path] = []
        self.merge_output_folder = tk.StringVar()
        self.merge_file_name = tk.StringVar()
        self.merge_password = tk.StringVar()

        self.split_pdf_path = tk.StringVar()
        self.split_output = tk.StringVar()
        self.split_pages = tk.StringVar()
        self.split_password = tk.StringVar()
        self.split_mode_label = tk.StringVar(value=self.tr("split_each_page"))

        self.office_pdf_path = tk.StringVar()
        self.office_pages = tk.StringVar()
        self.office_password = tk.StringVar()
        self.office_output_folder = tk.StringVar()
        self.office_word_file_name = tk.StringVar()
        self.office_ppt_file_name = tk.StringVar()

        self.source_office = tk.StringVar()
        self.office_pdf_output_folder = tk.StringVar()
        self.office_pdf_file_name = tk.StringVar()

        self.decrypt_source = tk.StringVar()
        self.decrypt_output_folder = tk.StringVar()
        self.decrypt_file_name = tk.StringVar()
        self.decrypt_password = tk.StringVar()

        self._build_ui()
        self.language_label.trace_add("write", self.on_language_change)

    def tr(self, key: str, **params: object) -> str:
        return translate(self.language, key, **params)

    def set_window_icon(self) -> None:
        icon_path = resource_path(APP_ICON_RESOURCE)
        if icon_path.is_file():
            try:
                self.iconbitmap(default=str(icon_path))
            except tk.TclError:
                pass

    def set_status(self, key: str, **params: object) -> None:
        self.status_key = key
        self.status_params = params
        self.status.set(self.tr(key, **params))

    def register_text(self, key: str, widget: tk.Widget) -> None:
        self.text_widgets.setdefault(key, []).append(widget)
        widget.configure(text=self.tr(key))

    def _build_ui(self) -> None:
        outer = ttk.Frame(self, padding=14)
        outer.pack(fill=tk.BOTH, expand=True)
        outer.columnconfigure(0, weight=1)
        outer.rowconfigure(1, weight=1)

        header = ttk.Frame(outer)
        header.grid(row=0, column=0, sticky="ew", pady=(0, 12))
        header.columnconfigure(0, weight=1)

        title = ttk.Label(header, font=("Segoe UI", 18, "bold"))
        title.grid(row=0, column=0, sticky="w")
        self.register_text("window_title", title)

        language_row = ttk.Frame(header)
        language_row.grid(row=0, column=1, sticky="e")
        language_label = ttk.Label(language_row)
        language_label.pack(side=tk.LEFT, padx=(0, 8))
        self.register_text("language", language_label)
        self.language_box = ttk.Combobox(
            language_row,
            textvariable=self.language_label,
            values=list(LANGUAGE_LABELS.values()),
            state="readonly",
            width=10,
        )
        self.language_box.pack(side=tk.LEFT)

        self.notebook = ttk.Notebook(outer)
        self.notebook.grid(row=1, column=0, sticky="nsew")

        self.build_images_tab()
        self.build_merge_tab()
        self.build_split_tab()
        self.build_pdf_to_office_tab()
        self.build_office_to_pdf_tab()
        self.build_decrypt_tab()

        status_row = ttk.Frame(outer)
        status_row.grid(row=2, column=0, sticky="ew", pady=(12, 0))
        status_row.columnconfigure(0, weight=1)
        ttk.Progressbar(status_row, maximum=100, variable=self.progress).grid(row=0, column=0, sticky="ew", padx=(0, 12))
        ttk.Label(status_row, textvariable=self.status, width=45).grid(row=0, column=1, sticky="e")

    def make_tab(self, key: str) -> ttk.Frame:
        frame = ttk.Frame(self.notebook, padding=16)
        frame.columnconfigure(1, weight=1)
        self.notebook.add(frame, text=self.tr(key))
        self.tabs[key] = frame
        return frame

    def labeled_entry(
        self,
        parent: ttk.Frame,
        row: int,
        label_key: str,
        variable: tk.StringVar,
        browse_command=None,
        show: str | None = None,
    ) -> ttk.Entry:
        label = ttk.Label(parent)
        label.grid(row=row, column=0, sticky="w", pady=7)
        self.register_text(label_key, label)
        entry = ttk.Entry(parent, textvariable=variable, show=show)
        entry.grid(row=row, column=1, sticky="ew", padx=8, pady=7)
        if browse_command:
            button = ttk.Button(parent, command=browse_command)
            button.grid(row=row, column=2, sticky="ew", pady=7)
            self.register_text("browse", button)
        return entry

    def add_action_button(self, parent: ttk.Frame, row: int, key: str, command) -> ttk.Button:
        button = ttk.Button(parent, command=command)
        button.grid(row=row, column=2, sticky="ew", pady=(16, 7))
        self.register_text(key, button)
        return button

    def build_images_tab(self) -> None:
        tab = self.make_tab("tab_images")
        self.labeled_entry(tab, 0, "pdf_file", self.image_pdf, lambda: self.pick_pdf(self.image_pdf))
        self.labeled_entry(tab, 1, "output_folder", self.image_output, lambda: self.pick_folder(self.image_output))
        self.labeled_entry(tab, 2, "password", self.image_password, show="*")
        self.add_hint(tab, 2, "password_hint")
        self.add_dpi_row(tab, 3, self.image_dpi)
        self.labeled_entry(tab, 4, "pages", self.image_pages)
        self.add_hint(tab, 4, "pages_hint")
        self.add_action_button(tab, 5, "export_images", self.start_export_images)

    def build_merge_tab(self) -> None:
        tab = self.make_tab("tab_merge")
        label = ttk.Label(tab)
        label.grid(row=0, column=0, sticky="nw", pady=7)
        self.register_text("pdf_files_to_merge", label)

        self.merge_listbox = tk.Listbox(tab, height=9)
        self.merge_listbox.grid(row=0, column=1, sticky="nsew", padx=8, pady=7)
        tab.rowconfigure(0, weight=1)

        buttons = ttk.Frame(tab)
        buttons.grid(row=0, column=2, sticky="nsew", pady=7)
        for key, command in [
            ("add_files", self.add_merge_files),
            ("remove_selected", self.remove_selected_merge_file),
            ("move_up", lambda: self.move_merge_file(-1)),
            ("move_down", lambda: self.move_merge_file(1)),
            ("clear", self.clear_merge_files),
        ]:
            button = ttk.Button(buttons, command=command)
            button.pack(fill=tk.X, pady=3)
            self.register_text(key, button)

        self.labeled_entry(tab, 1, "merge_output_folder", self.merge_output_folder, lambda: self.pick_folder(self.merge_output_folder))
        self.labeled_entry(tab, 2, "merge_file_name", self.merge_file_name)
        self.add_hint(tab, 2, "merge_name_hint")
        self.labeled_entry(tab, 3, "password", self.merge_password, show="*")
        self.add_hint(tab, 3, "password_hint")
        self.add_action_button(tab, 4, "merge", self.start_merge)

    def build_split_tab(self) -> None:
        tab = self.make_tab("tab_split")
        self.labeled_entry(tab, 0, "pdf_file", self.split_pdf_path, lambda: self.pick_pdf(self.split_pdf_path))
        self.labeled_entry(tab, 1, "output_folder", self.split_output, lambda: self.pick_folder(self.split_output))
        self.labeled_entry(tab, 2, "password", self.split_password, show="*")
        self.add_hint(tab, 2, "password_hint")
        self.labeled_entry(tab, 3, "pages", self.split_pages)
        self.add_hint(tab, 3, "pages_hint")

        label = ttk.Label(tab)
        label.grid(row=4, column=0, sticky="w", pady=7)
        self.register_text("split_mode", label)
        self.split_mode_box = ttk.Combobox(
            tab,
            textvariable=self.split_mode_label,
            values=[self.tr("split_each_page"), self.tr("split_one_file")],
            state="readonly",
        )
        self.split_mode_box.grid(row=4, column=1, sticky="w", padx=8, pady=7)
        self.add_action_button(tab, 5, "split", self.start_split)

    def build_pdf_to_office_tab(self) -> None:
        tab = self.make_tab("tab_pdf_to_office")
        self.labeled_entry(tab, 0, "pdf_file", self.office_pdf_path, lambda: self.pick_pdf(self.office_pdf_path))
        self.labeled_entry(tab, 1, "password", self.office_password, show="*")
        self.add_hint(tab, 1, "password_hint")
        self.labeled_entry(tab, 2, "pages", self.office_pages)
        self.add_hint(tab, 2, "pages_hint")
        self.labeled_entry(tab, 3, "merge_output_folder", self.office_output_folder, lambda: self.pick_folder(self.office_output_folder))
        self.labeled_entry(tab, 4, "word_file_name", self.office_word_file_name)
        self.add_hint(tab, 4, "word_name_hint")
        self.labeled_entry(tab, 5, "ppt_file_name", self.office_ppt_file_name)
        self.add_hint(tab, 5, "ppt_name_hint")

        button_row = ttk.Frame(tab)
        button_row.grid(row=6, column=1, columnspan=2, sticky="e", pady=(16, 7))
        word_button = ttk.Button(button_row, command=self.start_pdf_to_word)
        word_button.pack(side=tk.LEFT, padx=4)
        self.register_text("to_word", word_button)
        ppt_button = ttk.Button(button_row, command=self.start_pdf_to_ppt)
        ppt_button.pack(side=tk.LEFT, padx=4)
        self.register_text("to_ppt", ppt_button)

    def build_office_to_pdf_tab(self) -> None:
        tab = self.make_tab("tab_office_to_pdf")
        self.labeled_entry(tab, 0, "office_file", self.source_office, lambda: self.pick_office(self.source_office))
        self.labeled_entry(tab, 1, "merge_output_folder", self.office_pdf_output_folder, lambda: self.pick_folder(self.office_pdf_output_folder))
        self.labeled_entry(tab, 2, "pdf_file_name", self.office_pdf_file_name)
        self.add_hint(tab, 2, "office_pdf_name_hint")
        self.add_action_button(tab, 3, "to_pdf", self.start_office_to_pdf)

    def build_decrypt_tab(self) -> None:
        tab = self.make_tab("tab_decrypt")
        self.labeled_entry(tab, 0, "pdf_file", self.decrypt_source, lambda: self.pick_pdf(self.decrypt_source))
        self.labeled_entry(tab, 1, "password", self.decrypt_password, show="*")
        self.labeled_entry(tab, 2, "merge_output_folder", self.decrypt_output_folder, lambda: self.pick_folder(self.decrypt_output_folder))
        self.labeled_entry(tab, 3, "pdf_file_name", self.decrypt_file_name)
        self.add_hint(tab, 3, "decrypt_pdf_name_hint")
        self.add_action_button(tab, 4, "decrypt", self.start_decrypt)

    def add_hint(self, parent: ttk.Frame, row: int, key: str) -> None:
        hint = ttk.Label(parent, foreground="#666666")
        hint.grid(row=row, column=2, sticky="w", pady=7)
        self.register_text(key, hint)

    def add_dpi_row(self, parent: ttk.Frame, row: int, variable: tk.IntVar) -> None:
        label = ttk.Label(parent)
        label.grid(row=row, column=0, sticky="w", pady=7)
        self.register_text("dpi", label)
        spinner = ttk.Spinbox(parent, from_=72, to=600, increment=25, textvariable=variable, width=10)
        spinner.grid(row=row, column=1, sticky="w", padx=8, pady=7)

    def on_language_change(self, *_args: object) -> None:
        selected_language = LABEL_TO_LANGUAGE.get(self.language_label.get())
        if not selected_language or selected_language == self.language:
            return

        previous_split_mode = self.split_mode_label.get()
        was_one_file = previous_split_mode == self.tr("split_one_file")
        self.language = selected_language
        self.title(self.tr("app_title"))

        for key, widgets in self.text_widgets.items():
            for widget in widgets:
                widget.configure(text=self.tr(key))

        for key, frame in self.tabs.items():
            self.notebook.tab(frame, text=self.tr(key))

        self.split_mode_box.configure(values=[self.tr("split_each_page"), self.tr("split_one_file")])
        self.split_mode_label.set(self.tr("split_one_file" if was_one_file else "split_each_page"))
        self.status.set(self.tr(self.status_key, **self.status_params))

    def format_error(self, exc: Exception) -> str:
        if isinstance(exc, AppError):
            return self.tr(exc.message_key, **exc.params)
        return str(exc)

    def show_error(self, exc_or_key) -> None:
        text = self.format_error(exc_or_key) if isinstance(exc_or_key, Exception) else self.tr(exc_or_key)
        messagebox.showerror(self.tr("app_title"), text)

    def run_task(self, task, success_key: str, success_params) -> None:
        self.progress.set(0)
        self.set_status("running")

        def worker() -> None:
            try:
                result_params = task()
                params = success_params(result_params)
                self.after(0, self.finish_task, success_key, params)
            except Exception as exc:
                self.after(0, self.fail_task, exc)

        threading.Thread(target=worker, daemon=True).start()

    def finish_task(self, success_key: str, params: dict[str, object]) -> None:
        self.progress.set(100)
        detail = self.tr(success_key, **params)
        self.set_status("done", detail=detail)
        messagebox.showinfo(self.tr("app_title"), detail)

    def fail_task(self, exc: Exception) -> None:
        self.set_status("failed")
        self.progress.set(0)
        self.show_error(exc)

    def make_progress_callback(self):
        def callback(percent: int, page: int, current: int, total: int) -> None:
            self.after(
                0,
                lambda: (
                    self.progress.set(percent),
                    self.set_status("progress_page", page=page, current=current, total=total),
                ),
            )

        return callback

    def pick_pdf(self, variable: tk.StringVar) -> None:
        filename = filedialog.askopenfilename(
            title=self.tr("select_pdf_title"),
            filetypes=[(self.tr("pdf_files"), "*.pdf"), (self.tr("all_files"), "*.*")],
        )
        if filename:
            variable.set(filename)

    def pick_pdfs(self) -> list[Path]:
        filenames = filedialog.askopenfilenames(
            title=self.tr("select_pdfs_title"),
            filetypes=[(self.tr("pdf_files"), "*.pdf"), (self.tr("all_files"), "*.*")],
        )
        return [Path(name) for name in filenames]

    def pick_folder(self, variable: tk.StringVar) -> None:
        folder = filedialog.askdirectory(title=self.tr("select_output_title"))
        if folder:
            variable.set(folder)

    def pick_save_pdf(self, variable: tk.StringVar) -> None:
        filename = filedialog.asksaveasfilename(
            title=self.tr("select_output_pdf_title"),
            defaultextension=".pdf",
            filetypes=[(self.tr("pdf_files"), "*.pdf")],
        )
        if filename:
            variable.set(filename)

    def pick_save_docx(self, variable: tk.StringVar) -> None:
        filename = filedialog.asksaveasfilename(
            title=self.tr("select_output_docx_title"),
            defaultextension=".docx",
            filetypes=[(self.tr("word_files"), "*.docx")],
        )
        if filename:
            variable.set(filename)

    def pick_save_pptx(self, variable: tk.StringVar) -> None:
        filename = filedialog.asksaveasfilename(
            title=self.tr("select_output_pptx_title"),
            defaultextension=".pptx",
            filetypes=[(self.tr("ppt_files"), "*.pptx")],
        )
        if filename:
            variable.set(filename)

    def pick_office(self, variable: tk.StringVar) -> None:
        filename = filedialog.askopenfilename(
            title=self.tr("select_office_title"),
            filetypes=[
                (self.tr("office_files"), "*.doc *.docx *.ppt *.pptx"),
                (self.tr("all_files"), "*.*"),
            ],
        )
        if filename:
            variable.set(filename)

    def get_dpi(self, variable: tk.IntVar) -> int:
        try:
            dpi = int(variable.get())
        except (tk.TclError, ValueError) as exc:
            raise AppError("dpi_number") from exc
        if dpi < 72 or dpi > 600:
            raise AppError("dpi_range")
        return dpi

    def start_export_images(self) -> None:
        def task():
            pdf = Path(self.image_pdf.get())
            output = Path(self.image_output.get())
            ensure_pdf(pdf)
            ensure_output_folder(output)
            count = export_pdf_pages_to_images(
                pdf,
                output,
                self.get_dpi(self.image_dpi),
                self.image_pages.get(),
                self.image_password.get(),
                self.make_progress_callback(),
            )
            return {"count": count, "output": output}

        self.run_task(task, "images_done", lambda result: result)

    def add_merge_files(self) -> None:
        for path in self.pick_pdfs():
            if path not in self.merge_files:
                self.merge_files.append(path)
                self.merge_listbox.insert(tk.END, str(path))

    def remove_selected_merge_file(self) -> None:
        selected = list(self.merge_listbox.curselection())
        for index in reversed(selected):
            del self.merge_files[index]
            self.merge_listbox.delete(index)

    def clear_merge_files(self) -> None:
        self.merge_files.clear()
        self.merge_listbox.delete(0, tk.END)

    def move_merge_file(self, direction: int) -> None:
        selected = self.merge_listbox.curselection()
        if not selected:
            return
        index = selected[0]
        new_index = index + direction
        if new_index < 0 or new_index >= len(self.merge_files):
            return
        self.merge_files[index], self.merge_files[new_index] = self.merge_files[new_index], self.merge_files[index]
        text = self.merge_listbox.get(index)
        self.merge_listbox.delete(index)
        self.merge_listbox.insert(new_index, text)
        self.merge_listbox.selection_set(new_index)

    def start_merge(self) -> None:
        def task():
            if len(self.merge_files) < 2:
                raise AppError("need_two_pdfs")
            for pdf in self.merge_files:
                ensure_pdf(pdf)
            output = build_merge_output_pdf(
                self.merge_files,
                self.merge_output_folder.get(),
                self.merge_file_name.get(),
            )
            count = merge_pdfs(self.merge_files, output, self.merge_password.get())
            return {"count": count, "output": output}

        self.run_task(task, "merge_done", lambda result: result)

    def start_split(self) -> None:
        def task():
            pdf = Path(self.split_pdf_path.get())
            output = Path(self.split_output.get())
            ensure_pdf(pdf)
            ensure_output_folder(output)
            one_file = self.split_mode_label.get() == self.tr("split_one_file")
            count = split_pdf(pdf, output, self.split_pages.get(), self.split_password.get(), one_file)
            return {"count": count, "output": output}

        self.run_task(task, "split_done", lambda result: result)

    def start_pdf_to_word(self) -> None:
        def task():
            pdf = Path(self.office_pdf_path.get())
            ensure_pdf(pdf)
            output = build_named_output_file(
                pdf,
                self.office_output_folder.get(),
                self.office_word_file_name.get(),
                "PDF转Word",
                ".docx",
            )
            pdf_to_word(pdf, output, self.office_pages.get(), self.office_password.get(), self.make_progress_callback())
            return {"output": output}

        self.run_task(task, "office_done", lambda result: result)

    def start_pdf_to_ppt(self) -> None:
        def task():
            pdf = Path(self.office_pdf_path.get())
            ensure_pdf(pdf)
            output = build_named_output_file(
                pdf,
                self.office_output_folder.get(),
                self.office_ppt_file_name.get(),
                "PDF转PPT",
                ".pptx",
            )
            pdf_to_ppt(pdf, output, self.office_pages.get(), self.office_password.get(), self.make_progress_callback())
            return {"output": output}

        self.run_task(task, "office_done", lambda result: result)

    def start_office_to_pdf(self) -> None:
        def task():
            source = Path(self.source_office.get())
            if not source.is_file():
                raise AppError("valid_office")
            output = build_named_output_file(
                source,
                self.office_pdf_output_folder.get(),
                self.office_pdf_file_name.get(),
                office_to_pdf_default_prefix(source),
                ".pdf",
            )
            office_to_pdf(source, output)
            return {"output": output}

        self.run_task(task, "office_done", lambda result: result)

    def start_decrypt(self) -> None:
        def task():
            source = Path(self.decrypt_source.get())
            ensure_pdf(source)
            output = build_named_output_file(
                source,
                self.decrypt_output_folder.get(),
                self.decrypt_file_name.get(),
                "解密PDF",
                ".pdf",
            )
            decrypt_pdf(source, output, self.decrypt_password.get())
            return {"output": output}

        self.run_task(task, "decrypt_done", lambda result: result)


if __name__ == "__main__":
    app = PdfToolsApp()
    app.mainloop()
